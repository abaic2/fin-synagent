# -*- coding: utf-8 -*-
"""提取 PPTX / DOCX 全部文本与结构信息到 JSON（editor_sdk 失败时的后备解析器）

用法:
    python extract_office_content.py <输出.json> <输入文件1> [输入文件2 ...]

依赖: python-pptx, python-docx
"""
import sys
import json


def shape_text(shape, depth=0):
    out = []
    if shape.shape_type == 6:  # group
        for s in shape.shapes:
            out.extend(shape_text(s, depth + 1))
        return out
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            t = "".join(run.text for run in para.runs).strip()
            if t:
                sizes = [r.font.size.pt for r in para.runs if r.font.size]
                out.append({"text": t, "level": para.level,
                            "size": max(sizes) if sizes else None,
                            "bold": any(r.font.bold for r in para.runs)})
    if getattr(shape, "has_table", False) and shape.has_table:
        out.append({"table": [[c.text.strip() for c in r.cells] for r in shape.table.rows]})
    if shape.shape_type == 13:
        out.append({"picture": True})
    if getattr(shape, "has_chart", False) and shape.has_chart:
        try:
            ch = shape.chart
            out.append({"chart": {"type": str(ch.chart_type),
                                  "categories": list(ch.plots[0].categories),
                                  "series": [{"name": s.name, "values": list(s.values)} for s in ch.series]}})
        except Exception:
            out.append({"chart": {"type": "unknown"}})
    return out


def extract_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides):
        items = []
        for shape in slide.shapes:
            for item in shape_text(shape):
                item["x"] = round(shape.left / 12700, 1) if shape.left is not None else None
                item["y"] = round(shape.top / 12700, 1) if shape.top is not None else None
                items.append(item)
        notes = slide.notes_slide.notes_text_frame.text.strip() if slide.has_notes_slide else ""
        slides.append({"index": i, "items": items, "notes": notes})
    return {"type": "pptx", "slides": slides}


def extract_docx(path):
    from docx import Document
    doc = Document(path)
    blocks = []
    for para in doc.paragraphs:
        t = para.text.strip()
        if t:
            blocks.append({"style": para.style.name if para.style else None, "text": t})
    for tbl in doc.tables:
        blocks.append({"table": [[c.text.strip() for c in r.cells] for r in tbl.rows]})
    return {"type": "docx", "blocks": blocks}


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print(__doc__)
        sys.exit(1)
    out_path, inputs = sys.argv[1], sys.argv[2:]
    result = {}
    for p in inputs:
        low = p.lower()
        if low.endswith(".pptx") or low.endswith(".ppt"):
            result[p] = extract_pptx(p)
        elif low.endswith(".docx") or low.endswith(".doc"):
            result[p] = extract_docx(p)
        else:
            print(f"skip unsupported: {p}")
    with open(out_path, "w", encoding="utf-8") as f:
        json.dump(result, f, ensure_ascii=False, indent=1)
    for p, d in result.items():
        n = len(d.get("slides", d.get("blocks", [])))
        print(f"OK {p} -> {d['type']} {n} 个单元")
    print(f"saved -> {out_path}")
